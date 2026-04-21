#!/usr/bin/env python3
"""
PPD Lab — Course Enrollment System  ·  Presentation deck
Experiments 1-10 as per course syllabus:
  1  Problem Statement
  2  IEEE SRS + Risk Management + Gantt Chart
  3  Use Case Model + UML Activity Diagram
  4  UML Class Diagram + Interaction Diagrams
  5  State Chart Diagram
  6  Layered Architecture + UML Package Diagram
  7  Technical Services & Domain Objects Layer (implementation)
  8  User Interface Layer (implementation)
  9  Component & Deployment Diagrams
  10 Testing, Documentation & Final Product
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUT  = ROOT / "docs" / "PPD-LAB-Overview.pptx"

W = Inches(10)
H = Inches(7.5)

C = {
    "navy":   RGBColor(0x1A, 0x2B, 0x5E),
    "blue":   RGBColor(0x23, 0x6F, 0xB4),
    "sky":    RGBColor(0xD9, 0xEA, 0xF9),
    "teal":   RGBColor(0x2A, 0x9D, 0x8F),
    "lteal":  RGBColor(0xD0, 0xF0, 0xEB),
    "orange": RGBColor(0xE7, 0x6F, 0x51),
    "lorn":   RGBColor(0xFD, 0xE8, 0xE2),
    "green":  RGBColor(0x2D, 0x6A, 0x4F),
    "lgrn":   RGBColor(0xD8, 0xF3, 0xDC),
    "purple": RGBColor(0x6A, 0x3C, 0x9A),
    "lpurp":  RGBColor(0xEF, 0xE0, 0xFF),
    "yellow": RGBColor(0xF4, 0xA2, 0x61),
    "lyell":  RGBColor(0xFF, 0xF3, 0xCD),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
    "dark":   RGBColor(0x1A, 0x1A, 0x2E),
    "mid":    RGBColor(0x55, 0x5B, 0x6E),
    "light":  RGBColor(0xF0, 0xF4, 0xF8),
    "border": RGBColor(0xCC, 0xD6, 0xE8),
}


# ─────────────────────────── primitives ─────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0), radius=False):
    geom = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp  = slide.shapes.add_shape(geom, l, t, w, h)
    if fill:  shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:     shp.fill.background()
    if border: shp.line.color.rgb = border; shp.line.width = bw
    else:      shp.line.fill.background()
    return shp

def _ellipse(slide, l, t, w, h, fill=None, border=None, bw=Pt(1)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    if fill:  shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:     shp.fill.background()
    if border: shp.line.color.rgb = border; shp.line.width = bw
    else:      shp.line.fill.background()
    return shp

def _diamond(slide, l, t, w, h, fill=C["lyell"], border=C["yellow"]):
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border; shp.line.width = Pt(1.2)
    return shp

def _text(slide, l, t, w, h, txt, size=Pt(13), bold=False, color=None,
          align=PP_ALIGN.LEFT, italic=False, wrap=True, va=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = va
    p  = tf.paragraphs[0]; p.text = txt; p.alignment = align
    r  = p.runs[0] if p.runs else p.add_run()
    r.font.size = size; r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = color
    return tb

def _stext(shp, txt, size=Pt(13), bold=False, color=None,
           align=PP_ALIGN.CENTER, wrap=True):
    tf = shp.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p  = tf.paragraphs[0]; p.text = txt; p.alignment = align
    r  = p.runs[0] if p.runs else p.add_run()
    r.font.size = size; r.font.bold = bold
    if color: r.font.color.rgb = color

def _line(slide, x1, y1, x2, y2, color=C["mid"], w=Pt(1.5)):
    c = slide.shapes.add_connector(1, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = w
    return c


# ─────────────────────────── composite helpers ───────────────────────────────

def _header(slide, title, sub=""):
    _rect(slide, 0, 0, W, Inches(1.0), fill=C["navy"])
    _text(slide, Inches(0.4), Inches(0.1), Inches(9.2), Inches(0.6),
          title, size=Pt(24), bold=True, color=C["white"], va=MSO_ANCHOR.MIDDLE)
    if sub:
        _text(slide, Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.32),
              sub, size=Pt(12), italic=True,
              color=RGBColor(0xB0, 0xC4, 0xDE))

def _footer(slide, txt="PPD Lab — Course Enrollment System"):
    _rect(slide, 0, Inches(7.18), W, Inches(0.32), fill=C["navy"])
    _text(slide, Inches(0.3), Inches(7.2), Inches(9.4), Inches(0.28),
          txt, size=Pt(10), color=RGBColor(0xB0, 0xC4, 0xDE),
          align=PP_ALIGN.CENTER)

def _divider(prs, exp, title, sub):
    slide = _blank(prs)
    _rect(slide, 0, 0, W, H, fill=C["navy"])
    _rect(slide, 0, Inches(3.8), W, Inches(0.15), fill=C["teal"])
    _text(slide, 0, Inches(0.7), W, Inches(1.4),
          exp, size=Pt(60), bold=True, color=RGBColor(0x3A,0x5A,0xAA),
          align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.6),
          title, size=Pt(28), bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.2), Inches(4.1), Inches(7.6), Inches(0.9),
          sub, size=Pt(15), italic=True,
          color=RGBColor(0xC0, 0xD8, 0xF8), align=PP_ALIGN.CENTER)

def _boxt(slide, l, t, w, h, txt, fill=C["sky"], border=C["blue"],
          bw=Pt(1.2), size=Pt(12), bold=False, radius=True, tc=None):
    shp = _rect(slide, l, t, w, h, fill=fill, border=border, bw=bw, radius=radius)
    _stext(shp, txt, size=size, bold=bold, color=tc or C["dark"])
    return shp

def _card(slide, l, t, w, h, heading, items, hfill=C["blue"], bfill=C["sky"]):
    _rect(slide, l, t, w, Inches(0.4), fill=hfill, radius=False)
    _text(slide, l+Inches(0.1), t+Inches(0.02), w-Inches(0.2), Inches(0.38),
          heading, size=Pt(12), bold=True, color=C["white"], va=MSO_ANCHOR.MIDDLE)
    bh = h - Inches(0.4)
    _rect(slide, l, t+Inches(0.4), w, bh, fill=bfill,
          border=C["border"], bw=Pt(0.8), radius=False)
    y = t + Inches(0.5)
    for item in items:
        _text(slide, l+Inches(0.12), y, w-Inches(0.2), Inches(0.32),
              f"• {item}", size=Pt(11), color=C["dark"])
        y += Inches(0.29)

def _uml_class(slide, l, t, w, name, attrs, methods,
               hfill=C["navy"], bfill=C["sky"]):
    rh = Inches(0.38)
    _rect(slide, l, t, w, rh, fill=hfill, radius=False)
    _text(slide, l+Inches(0.05), t+Inches(0.02), w-Inches(0.1), rh-Inches(0.04),
          name, size=Pt(12), bold=True, color=C["white"],
          align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    y = t + rh
    for a in attrs:
        _rect(slide, l, y, w, rh, fill=bfill, border=C["border"], bw=Pt(0.4))
        _text(slide, l+Inches(0.08), y+Inches(0.02), w-Inches(0.12), rh,
              a, size=Pt(10), color=C["dark"], va=MSO_ANCHOR.MIDDLE)
        y += rh
    _rect(slide, l, y, w, Inches(0.02), fill=C["border"])
    y += Inches(0.02)
    for m in methods:
        _rect(slide, l, y, w, rh, fill=C["white"], border=C["border"], bw=Pt(0.4))
        _text(slide, l+Inches(0.08), y+Inches(0.02), w-Inches(0.12), rh,
              m, size=Pt(10), color=C["mid"], italic=True, va=MSO_ANCHOR.MIDDLE)
        y += rh
    return y  # bottom edge


# ─────────────────────────── individual slides ───────────────────────────────

def s_title(prs):
    slide = _blank(prs)
    _rect(slide, 0, 0, W, H, fill=C["navy"])
    _rect(slide, 0, Inches(5.1), W, Inches(0.14), fill=C["teal"])
    _text(slide, Inches(0.7), Inches(1.4), Inches(8.6), Inches(1.5),
          "PPD Lab\nCourse Enrollment System",
          size=Pt(36), bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    _text(slide, Inches(0.7), Inches(3.1), Inches(8.6), Inches(0.8),
          "A full-stack student–admin platform built with React + Fastify + PostgreSQL",
          size=Pt(16), italic=True,
          color=RGBColor(0xB0,0xC4,0xDE), align=PP_ALIGN.CENTER)
    tags = ["Exp 1: Problem", "Exp 2: SRS+Gantt", "Exp 3: Use Case+Activity",
            "Exp 4: Class+Interaction", "Exp 5: State Chart"]
    tags2 = ["Exp 6: Package Arch", "Exp 7: Tech+Domain Impl",
             "Exp 8: UI Layer", "Exp 9: Comp+Deploy", "Exp 10: Test+Docs"]
    for row, tlist in enumerate([tags, tags2]):
        for i, t in enumerate(tlist):
            x = Inches(0.35) + i * Inches(1.88)
            y = Inches(5.38) + row * Inches(0.68)
            _boxt(slide, x, y, Inches(1.82), Inches(0.55), t,
                  fill=RGBColor(0x26,0x3A,0x70),
                  border=C["teal"], size=Pt(9.5), tc=C["white"])
    _text(slide, 0, Inches(7.1), W, Inches(0.38),
          "PPD Lab Coursework  ·  2026",
          size=Pt(11), color=RGBColor(0x70,0x88,0xAA), align=PP_ALIGN.CENTER)


def s_intro(prs):
    slide = _blank(prs)
    _header(slide, "Introduction", "Project overview and context")
    _footer(slide)
    _card(slide, Inches(0.3), Inches(1.1), Inches(4.5), Inches(5.9),
          "About the Project",
          ["Web-based course enrollment and student administration",
           "Monorepo: 3 React/Vite SPAs + 1 Fastify REST API",
           "PostgreSQL (Neon) via Drizzle ORM",
           "JWT authentication + RBAC (student / admin)",
           "Follows 10-experiment PPD Lab syllabus",
           "Phase 1→10: problem → SRS → design → implementation → test",
           "OpenAPI/Swagger UI at /documentation",
           "GitHub Actions CI: lint, test, build on every push"],
          hfill=C["navy"])
    _card(slide, Inches(5.0), Inches(1.1), Inches(4.65), Inches(2.75),
          "Technology Stack",
          ["React 19, Vite 8, TypeScript, shadcn/ui (Base UI), Tailwind",
           "Node.js 20+, Fastify, Drizzle ORM, Zod validation",
           "PostgreSQL on Neon (pooled connection string)",
           "JWT (jose) + bcrypt  ·  Pino structured logging",
           "Vitest (unit + HTTP smoke)  ·  GitHub Actions CI"],
          hfill=C["blue"])
    _card(slide, Inches(5.0), Inches(4.05), Inches(4.65), Inches(2.95),
          "User Roles",
          ["Student: register, login, browse catalog, request & cancel enrollment,",
           "  view status, update profile, receive notifications",
           "Admin: manage courses & students, approve/reject enrollments,",
           "  view dashboard & reports, manage admin accounts,",
           "  send notifications to students"],
          hfill=C["teal"])


def s_objective(prs):
    slide = _blank(prs)
    _header(slide, "Objective", "Goals and expected outcomes of this project")
    _footer(slide)
    goals = [
        ("Primary Goal",
         ["Build an end-to-end course enrollment platform for a university context",
          "Two roles: Student (self-service) and Admin (governance)",
          "Persistent, consistent data via PostgreSQL with Drizzle migrations"],
         C["navy"]),
        ("Security & Authorization",
         ["Passwords stored as bcrypt hashes; JWT tokens with per-token jti",
          "Role-based access control on every route (student vs admin)",
          "Token revocation on logout via revoked_tokens table"],
         C["blue"]),
        ("Business Rules",
         ["Capacity enforcement: approved_count < capacity before approve",
          "Prerequisite check at enrollment request time",
          "Duplicate enrollment prevention (UNIQUE constraint)"],
         C["teal"]),
        ("Engineering Quality",
         ["TypeScript end-to-end; Zod schemas for all inputs/outputs",
          "Vitest tests + optional DB integration tests",
          "CI pipeline (lint, test, build) on GitHub Actions"],
         C["green"]),
        ("User Experience",
         ["Clean React SPA with role-aware routing and navigation",
          "In-app notifications with bell + Sonner toasts",
          "Enrollment status polling every 15 s while tab visible"],
         C["purple"]),
        ("Documentation",
         ["OpenAPI/Swagger at /documentation auto-generated",
          "README with setup, env vars, migrate, first-admin steps",
          "This presentation: all 10 experiments covered"],
         C["orange"]),
    ]
    for i, (t, items, col) in enumerate(goals):
        row, c = divmod(i, 2)
        _card(slide, Inches(0.3) + c*Inches(4.85),
              Inches(1.1) + row*Inches(2.1),
              Inches(4.62), Inches(1.95), t, items, hfill=col)


# ── Experiment 1 ─────────────────────────────────────────────────────────────

def s_exp1(prs):
    _divider(prs, "Experiment 1", "Problem Statement",
             "Articulating the core challenges this system must address")

    slide = _blank(prs)
    _header(slide, "Experiment 1 — Problem Statement",
            "Why a web-based enrollment system is needed")
    _footer(slide)

    _rect(slide, Inches(0.3), Inches(1.08), Inches(9.4), Inches(0.7),
          fill=C["lorn"], border=C["orange"], bw=Pt(1.5), radius=True)
    _text(slide, Inches(0.5), Inches(1.12), Inches(9.0), Inches(0.64),
          "Manual / spreadsheet-based enrollment is error-prone, opaque, and difficult to audit. "
          "Students and administrators lack a single governed platform with real-time visibility, "
          "enforced rules, and secure role-separated access.",
          size=Pt(13), italic=True, color=RGBColor(0x70,0x20,0x00),
          wrap=True, va=MSO_ANCHOR.MIDDLE)

    problems = [
        ("P1 — Capacity Overruns",
         "Without server-side enforcement, courses can exceed seat limits. No automated feedback to students on availability.",
         C["orange"], C["lorn"]),
        ("P2 — Prerequisite Blindness",
         "Students may skip required prior courses. Manual checking is inconsistent and time-consuming for staff.",
         C["blue"], C["sky"]),
        ("P3 — No Role Separation",
         "Flat access models allow cross-role data leaks. Admin-only functions (approve, manage, report) must be strictly guarded.",
         C["purple"], C["lpurp"]),
        ("P4 — Weak Audit Trail",
         "Enrollment history, status changes, and admin actions must be persisted and queryable for reports and disputes.",
         C["teal"], C["lteal"]),
        ("P5 — Poor Feedback Loop",
         "Students have no real-time visibility into request status. A notification system is required to close this gap.",
         C["green"], C["lgrn"]),
        ("P6 — Maintainability Gap",
         "Codebase must be typed, validated, tested, and documented so any developer can confidently extend or deploy it.",
         C["navy"], C["light"]),
    ]
    for i, (title, desc, hfill, bfill) in enumerate(problems):
        col, row = divmod(i, 3)
        x = Inches(0.3) + col * Inches(4.85)
        y = Inches(2.0) + row * Inches(1.7)
        _rect(slide, x, y, Inches(4.62), Inches(0.38), fill=hfill)
        _text(slide, x+Inches(0.1), y+Inches(0.02), Inches(4.42), Inches(0.35),
              title, size=Pt(11.5), bold=True, color=C["white"], va=MSO_ANCHOR.MIDDLE)
        _rect(slide, x, y+Inches(0.38), Inches(4.62), Inches(1.18),
              fill=bfill, border=C["border"], bw=Pt(0.6))
        _text(slide, x+Inches(0.12), y+Inches(0.46), Inches(4.4), Inches(1.06),
              desc, size=Pt(11), color=C["dark"], wrap=True)


# ── Experiment 2 ─────────────────────────────────────────────────────────────

def s_exp2_srs(prs):
    _divider(prs, "Experiment 2", "IEEE SRS + Risk Management + Gantt Chart",
             "Structured requirements, risk table, and project timeline")

    slide = _blank(prs)
    _header(slide, "Experiment 2 — IEEE Standard SRS",
            "Software Requirements Specification (condensed)")
    _footer(slide)

    _card(slide, Inches(0.3), Inches(1.08), Inches(4.6), Inches(6.0),
          "Functional Requirements",
          ["FR-01  User registration (email + password)",
           "FR-02  Login / logout with JWT token",
           "FR-03  GET /auth/me — authenticated identity",
           "FR-04  Profile view and update (/account)",
           "FR-05  Browse course catalog (list + detail)",
           "FR-06  Request course enrollment",
           "FR-07  View own enrollment status",
           "FR-08  Cancel pending enrollment",
           "FR-09  Admin: CRUD courses + prerequisites",
           "FR-10  Admin: approve / reject enrollments",
           "FR-11  Admin: student list and profile edit",
           "FR-12  Admin: dashboard KPIs and reports",
           "FR-13  In-app notifications (create / read)",
           "FR-14  Admin: manage other admin accounts"],
          hfill=C["navy"])

    _card(slide, Inches(5.1), Inches(1.08), Inches(4.55), Inches(3.0),
          "Non-Functional Requirements",
          ["NFR-01  Passwords hashed (bcrypt); never plaintext",
           "NFR-02  JWT: signed, scoped role claims, jti revocation",
           "NFR-03  CORS scoped to CLIENT_ORIGIN env var",
           "NFR-04  GET /health + GET /health/db endpoints",
           "NFR-05  Structured logging via Pino (LOG_LEVEL)",
           "NFR-06  Zod validation on all request bodies",
           "NFR-07  OpenAPI / Swagger at /documentation",
           "NFR-08  Drizzle DB migrations (versioned schema)"],
          hfill=C["blue"])

    _card(slide, Inches(5.1), Inches(4.28), Inches(4.55), Inches(2.8),
          "System Constraints",
          ["Node.js >= 20.19.0  (Vite 8 requires it)",
           "PostgreSQL on Neon; pooled connection string required",
           "Multiple Vite origins comma-separated in CLIENT_ORIGIN",
           "File uploads stored locally (lab scope only)",
           "No WebSocket — polling used for real-time UI updates",
           "Browser: modern Chromium/Firefox"],
          hfill=C["teal"])


def s_exp2_risk_gantt(prs):
    slide = _blank(prs)
    _header(slide, "Experiment 2 — Risk Management & Gantt Chart",
            "Key project risks and phased timeline")
    _footer(slide)

    # Risk table
    _rect(slide, Inches(0.3), Inches(1.08), Inches(9.4), Inches(0.4),
          fill=C["navy"])
    for col, (label, wid, off) in enumerate([
        ("Risk", 2.2, 0.0), ("Likelihood", 1.5, 2.2),
        ("Impact", 1.4, 3.7), ("Mitigation", 4.1, 5.1)
    ]):
        _text(slide, Inches(0.4+off), Inches(1.1), Inches(wid), Inches(0.36),
              label, size=Pt(11.5), bold=True, color=C["white"],
              va=MSO_ANCHOR.MIDDLE)

    risks = [
        ("DB connection failure",  "Low",    "High",   "GET /health/db + Neon pooling"),
        ("JWT secret leak",        "Low",    "High",   "Env var; never in source code"),
        ("Capacity race condition","Medium", "Medium", "Atomic count check on approve"),
        ("Missing prerequisites",  "Medium", "Medium", "Server-side check at POST /enrollments"),
        ("CORS misconfiguration",  "Low",    "Medium", "CLIENT_ORIGIN env var; tested locally"),
        ("Test coverage gap",      "Medium", "Low",    "Vitest suite + RUN_DB_INTEGRATION flag"),
    ]
    for i, (risk, like, imp, mit) in enumerate(risks):
        bg = C["light"] if i % 2 == 0 else C["white"]
        lc = C["orange"] if like == "Medium" else (C["teal"] if like == "Low" else C["orange"])
        ic = C["orange"] if imp == "High" else (C["yellow"] if imp == "Medium" else C["teal"])
        y = Inches(1.48) + i * Inches(0.38)
        _rect(slide, Inches(0.3), y, Inches(9.4), Inches(0.38),
              fill=bg, border=C["border"], bw=Pt(0.5))
        _text(slide, Inches(0.4), y+Inches(0.02), Inches(2.1), Inches(0.35),
              risk, size=Pt(11), color=C["dark"], va=MSO_ANCHOR.MIDDLE)
        _boxt(slide, Inches(2.65), y+Inches(0.04), Inches(1.15), Inches(0.3),
              like, fill=C["lorn"] if like == "Medium" else C["lteal"],
              border=lc, bw=Pt(0.8), size=Pt(10), radius=True, tc=C["dark"])
        _boxt(slide, Inches(3.95), y+Inches(0.04), Inches(1.05), Inches(0.3),
              imp, fill=C["lorn"] if imp == "High" else (C["lyell"] if imp == "Medium" else C["lteal"]),
              border=ic, bw=Pt(0.8), size=Pt(10), radius=True, tc=C["dark"])
        _text(slide, Inches(5.2), y+Inches(0.02), Inches(4.4), Inches(0.35),
              mit, size=Pt(11), color=C["dark"], va=MSO_ANCHOR.MIDDLE)

    # Gantt chart
    _text(slide, Inches(0.3), Inches(3.88), Inches(4.0), Inches(0.3),
          "Project Gantt Chart (Phases 1–10)", size=Pt(13), bold=True, color=C["navy"])

    phases = [
        ("Ph 1  Foundation",        0, 1,  C["blue"]),
        ("Ph 2  Authentication",    1, 1,  C["teal"]),
        ("Ph 3  Courses (SRS→UC→CD)", 2, 1.5, C["purple"]),
        ("Ph 4  Enrollment workflow", 3.5, 1.5, C["orange"]),
        ("Ph 5  Admin aggregation",   5, 1,  C["green"]),
        ("Ph 6  Notifications",       6, 0.8, C["blue"]),
        ("Ph 7  Tech+Domain impl",    6.8, 0.8, C["navy"]),
        ("Ph 8  UI layer",            7.6, 0.7, C["teal"]),
        ("Ph 9  Comp+Deploy diagrams",8.3, 0.7, C["purple"]),
        ("Ph 10 Test + Docs",         9, 1,  C["orange"]),
    ]
    timeline_w = Inches(6.2)
    total_units = 10.0
    unit = timeline_w / total_units
    gx = Inches(2.85); gy = Inches(4.26); bar_h = Inches(0.28); gap = Inches(0.04)
    for i, (label, start, dur, color) in enumerate(phases):
        y = gy + i * (bar_h + gap)
        _text(slide, Inches(0.3), y, Inches(2.5), bar_h,
              label, size=Pt(9.5), color=C["dark"], va=MSO_ANCHOR.MIDDLE)
        _rect(slide, gx + start * unit, y, dur * unit, bar_h,
              fill=color, border=None)


# ── Experiment 3 ─────────────────────────────────────────────────────────────

def s_exp3_usecase(prs):
    _divider(prs, "Experiment 3",
             "Use Case Model + UML Activity Diagram",
             "Who interacts with the system and what they do")

    slide = _blank(prs)
    _header(slide, "Experiment 3(i) — Use Case Diagram",
            "Actors: Student, Admin  |  System boundary: Course Enrollment System")
    _footer(slide)

    # System boundary
    _rect(slide, Inches(2.5), Inches(1.08), Inches(5.2), Inches(6.15),
          fill=C["light"], border=C["navy"], bw=Pt(2))
    _text(slide, Inches(2.5), Inches(1.1), Inches(5.2), Inches(0.3),
          "Course Enrollment System",
          size=Pt(11), bold=True, color=C["navy"], align=PP_ALIGN.CENTER)

    # Use cases inside boundary
    uc_student = [
        (Inches(2.65), Inches(1.55), "Register / Login"),
        (Inches(2.65), Inches(2.25), "View Profile"),
        (Inches(2.65), Inches(2.95), "Browse Catalog"),
        (Inches(2.65), Inches(3.65), "Request Enrollment"),
        (Inches(2.65), Inches(4.35), "View Enrollment Status"),
        (Inches(2.65), Inches(5.05), "Cancel Enrollment"),
        (Inches(2.65), Inches(5.75), "View Notifications"),
    ]
    uc_admin = [
        (Inches(4.9), Inches(1.55), "Manage Courses"),
        (Inches(4.9), Inches(2.25), "Manage Students"),
        (Inches(4.9), Inches(2.95), "Approve Enrollment"),
        (Inches(4.9), Inches(3.65), "Reject Enrollment"),
        (Inches(4.9), Inches(4.35), "View Dashboard"),
        (Inches(4.9), Inches(5.05), "Generate Reports"),
        (Inches(4.9), Inches(5.75), "Send Notifications"),
    ]
    ucw = Inches(2.25); uch = Inches(0.46)
    for x, y, label in uc_student + uc_admin:
        shp = _ellipse(slide, x, y, ucw, uch, fill=C["sky"], border=C["blue"], bw=Pt(1))
        _stext(shp, label, size=Pt(10.5), bold=False, color=C["navy"])

    # Actor — Student (left)
    _boxt(slide, Inches(0.3), Inches(3.45), Inches(1.85), Inches(0.56),
          "👤  Student", fill=C["lteal"], border=C["teal"],
          bw=Pt(1.5), size=Pt(13), bold=True, tc=C["dark"])
    # Actor — Admin (right)
    _boxt(slide, Inches(7.85), Inches(3.45), Inches(1.85), Inches(0.56),
          "🛡  Admin", fill=C["lorn"], border=C["orange"],
          bw=Pt(1.5), size=Pt(13), bold=True, tc=C["dark"])

    # Association lines
    sx = Inches(2.15); ax = Inches(4.9)
    for _, y, _ in uc_student:
        _line(slide, Inches(2.15), y+Inches(0.23), Inches(2.5), y+Inches(0.23),
              color=C["teal"], w=Pt(1.2))
    _line(slide, Inches(0.3)+Inches(0.93), Inches(3.73), Inches(2.15), Inches(3.73),
          color=C["teal"], w=Pt(1.5))

    for _, y, _ in uc_admin:
        _line(slide, Inches(4.9)+ucw, y+Inches(0.23), Inches(7.85), y+Inches(0.23),
              color=C["orange"], w=Pt(1.2))
    _line(slide, Inches(7.15), Inches(3.73), Inches(7.85), Inches(3.73),
          color=C["orange"], w=Pt(1.5))

    # Legend note
    _text(slide, Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.22),
          "«extend» Notify Student applies to both Approve Enrollment and Reject Enrollment  "
          "·  Capacity + prerequisite checks are «include» sub-flows of Approve Enrollment",
          size=Pt(9.5), italic=True, color=C["mid"])


def s_exp3_activity(prs):
    slide = _blank(prs)
    _header(slide, "Experiment 3(ii) — UML Activity Diagram",
            "Enrollment request business activity flow")
    _footer(slide)

    # Swimlane headers
    lanes = [("Student", C["sky"], C["blue"]),
             ("System / API", C["lteal"], C["teal"]),
             ("Admin", C["lorn"], C["orange"])]
    lw = Inches(3.0)
    for i, (name, fill, border) in enumerate(lanes):
        x = Inches(0.3) + i * lw
        _rect(slide, x, Inches(1.05), lw, Inches(0.38), fill=border)
        _text(slide, x+Inches(0.1), Inches(1.07), lw-Inches(0.2), Inches(0.35),
              name, size=Pt(13), bold=True, color=C["white"],
              align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
        _rect(slide, x, Inches(1.43), lw, Inches(5.65),
              fill=fill, border=border, bw=Pt(0.8))
        if i > 0:
            _rect(slide, x, Inches(1.43), Inches(0.02), Inches(5.65), fill=border)

    bw = Inches(2.6); bh = Inches(0.5)

    def act(lane, y, label, fill=C["white"], border=C["navy"]):
        x = Inches(0.5) + lane * Inches(3.0)
        shp = _rect(slide, x, y, bw, bh, fill=fill, border=border, bw=Pt(1.2), radius=True)
        _stext(shp, label, size=Pt(11), color=C["dark"])
        return x + bw/2, y + bh  # center-bottom

    def dec(lane, y, label):
        x = Inches(0.5) + lane * Inches(3.0) + Inches(0.3)
        shp = _diamond(slide, x, y, Inches(2.0), Inches(0.52))
        _stext(shp, label, size=Pt(10), color=C["dark"])
        return x + Inches(1.0), y + Inches(0.52)

    # Start
    _ellipse(slide, Inches(1.2), Inches(1.52), Inches(0.3), Inches(0.3),
             fill=C["navy"], border=C["navy"])
    _line(slide, Inches(1.35), Inches(1.82), Inches(1.35), Inches(1.95), color=C["navy"])

    cx, by = act(0, Inches(1.95), "Browse course catalog")
    _line(slide, cx, by, cx, by+Inches(0.12), color=C["navy"])
    cx2, by2 = act(0, by+Inches(0.12), "Click 'Request Enrollment'")
    _line(slide, cx2, by2, Inches(4.5), by2+Inches(0.12), color=C["navy"])

    # System checks
    cx3, by3 = act(1, by2+Inches(0.12), "Validate JWT + role", fill=C["lteal"], border=C["teal"])
    _line(slide, cx3, by3, cx3, by3+Inches(0.12), color=C["teal"])
    dx, dy = dec(1, by3+Inches(0.12), "Prereqs met?")
    _line(slide, dx, dy, dx, dy+Inches(0.12), color=C["teal"])
    cx4, by4 = act(1, dy+Inches(0.12), "Check capacity + duplicates", fill=C["lteal"], border=C["teal"])
    _line(slide, cx4, by4, cx4, by4+Inches(0.12), color=C["teal"])
    cx5, by5 = act(1, by4+Inches(0.12), "INSERT enrollment (PENDING)", fill=C["lteal"], border=C["teal"])

    # Admin
    _line(slide, cx5, by5, Inches(7.8), by5+Inches(0.12), color=C["orange"])
    cx6, by6 = act(2, by5+Inches(0.12), "Review pending enrollments", fill=C["lorn"], border=C["orange"])
    _line(slide, cx6, by6, cx6, by6+Inches(0.12), color=C["orange"])
    dxa, dya = dec(2, by6+Inches(0.12), "Approve?")
    _line(slide, dxa, dya, dxa, dya+Inches(0.12), color=C["orange"])
    act(2, dya+Inches(0.12), "UPDATE status → APPROVED/REJECTED\n+ notify student", fill=C["lorn"], border=C["orange"])

    # No path labels
    _text(slide, Inches(3.25), dy+Inches(0.18), Inches(0.8), Inches(0.28),
          "No → 422", size=Pt(9), italic=True, color=C["orange"])
    _text(slide, Inches(8.6), dya+Inches(0.18), Inches(0.8), Inches(0.28),
          "No →\nReject", size=Pt(9), italic=True, color=C["orange"])


# ── Experiment 4 ─────────────────────────────────────────────────────────────

def s_exp4_class(prs):
    _divider(prs, "Experiment 4",
             "UML Class Diagram + Interaction Diagrams",
             "Domain model & enrollment sequence")

    slide = _blank(prs)
    _header(slide, "Experiment 4(i) — UML Class Diagram",
            "Domain model: User, Course, Enrollment, Notification, CoursePrerequisite")
    _footer(slide)

    cw = Inches(2.5)
    # User class
    _uml_class(slide, Inches(0.25), Inches(1.1), cw, "User",
               ["+id: int (PK)", "+name: string", "+email: string",
                "+passwordHash: string", "+role: student|admin", "+isActive: bool"],
               ["+register()", "+login()", "+logout()", "+updateProfile()"],
               hfill=C["navy"])

    # Course class
    _uml_class(slide, Inches(3.85), Inches(1.1), cw, "Course",
               ["+id: int (PK)", "+code: string", "+title: string",
                "+description: string", "+credits: int", "+capacity: int"],
               ["+list()", "+detail()", "+create()", "+update()", "+delete()"],
               hfill=C["blue"])

    # Enrollment class
    _uml_class(slide, Inches(7.25), Inches(1.1), cw, "Enrollment",
               ["+id: int (PK)", "+userId: int (FK)", "+courseId: int (FK)",
                "+status: PENDING|APPROVED|REJECTED|CANCELLED",
                "+createdAt: timestamp"],
               ["+request()", "+approve()", "+reject()", "+cancel()"],
               hfill=C["teal"])

    # Notification
    _uml_class(slide, Inches(0.25), Inches(4.5), Inches(2.2), "Notification",
               ["+id: int (PK)", "+userId: int (FK)",
                "+message: string", "+read: bool"],
               ["+create()", "+markRead()"],
               hfill=C["orange"])

    # CoursePrerequisite (junction)
    _uml_class(slide, Inches(3.85), Inches(4.5), Inches(2.5), "CoursePrerequisite",
               ["<<junction>>", "+courseId: int (FK)",
                "+prerequisiteId: int (FK)"],
               ["+add()", "+remove()"],
               hfill=C["purple"])

    # RevokedToken
    _uml_class(slide, Inches(7.25), Inches(4.5), Inches(2.45), "RevokedToken",
               ["+id: int (PK)", "+jti: string",
                "+expiresAt: timestamp"],
               ["+revoke()", "+isRevoked()"],
               hfill=C["mid"])

    # Relationships
    # User 1—* Enrollment
    _line(slide, Inches(2.75), Inches(2.3), Inches(7.25), Inches(2.3),
          color=C["navy"], w=Pt(1.4))
    _text(slide, Inches(4.6), Inches(2.0), Inches(1.5), Inches(0.28),
          "1        *", size=Pt(11), bold=True, color=C["navy"])
    _text(slide, Inches(4.6), Inches(2.3), Inches(2.0), Inches(0.25),
          "has", size=Pt(10), italic=True, color=C["mid"], align=PP_ALIGN.CENTER)

    # Course 1—* Enrollment
    _line(slide, Inches(6.35), Inches(2.3), Inches(7.25)+Inches(0.8), Inches(2.3),
          color=C["blue"], w=Pt(1.4))
    # User 1—* Notification
    _line(slide, Inches(1.35), Inches(4.15), Inches(1.35), Inches(4.5),
          color=C["orange"], w=Pt(1.4))
    _text(slide, Inches(1.45), Inches(4.22), Inches(0.5), Inches(0.25),
          "1..*", size=Pt(10), color=C["orange"])
    # Course self-ref prerequisite
    _line(slide, Inches(4.85)+Inches(0.6), Inches(4.5),
          Inches(4.85)+Inches(0.6), Inches(4.18),
          color=C["purple"], w=Pt(1.4))
    _text(slide, Inches(5.2), Inches(4.25), Inches(1.5), Inches(0.25),
          "M:N prereq", size=Pt(10), italic=True, color=C["purple"])


def s_exp4_seq(prs):
    slide = _blank(prs)
    _header(slide, "Experiment 4(ii) — UML Interaction (Sequence) Diagram",
            "Scenario: student requests enrollment — objects and messages")
    _footer(slide)

    actors = ["StudentUI", "API\n/enrollments", "EnrollmentService", "DrizzleORM", "PostgreSQL"]
    xs = [Inches(0.6), Inches(2.6), Inches(4.6), Inches(6.6), Inches(8.5)]
    bw = Inches(1.55); bh = Inches(0.58)
    for x, a in zip(xs, actors):
        _boxt(slide, x, Inches(1.1), bw, bh, a, fill=C["navy"],
              border=C["teal"], tc=C["white"], size=Pt(11), bold=True)
        _rect(slide, x+Inches(0.74), Inches(1.68), Inches(0.06), Inches(5.35),
              fill=RGBColor(0xC0,0xD0,0xE8))

    def msg(y, fi, ti, label, ret=False, color=None):
        x1 = xs[fi]+Inches(0.77); x2 = xs[ti]+Inches(0.77)
        lx, rx = (x1, x2) if x2 > x1 else (x2, x1)
        c = slide.shapes.add_connector(1, lx, Inches(y), rx, Inches(y))
        col = color or (C["mid"] if ret else C["blue"])
        c.line.color.rgb = col; c.line.width = Pt(1.4)
        if ret: c.line.dash_style = 4
        mid = (lx + rx) / 2
        _text(slide, mid-Inches(1.1), Inches(y)-Inches(0.25), Inches(2.2), Inches(0.26),
              label, size=Pt(9.5), color=col, align=PP_ALIGN.CENTER, italic=ret)

    msg(2.05, 0, 1, "POST /enrollments {courseId}")
    msg(2.4,  1, 2, "checkPrerequisites(userId, courseId)", color=C["teal"])
    msg(2.72, 2, 3, "findApprovedEnrollments(userId)", color=C["teal"])
    msg(2.98, 3, 4, "SELECT enrollments WHERE user_id=…", color=C["mid"])
    msg(3.28, 4, 3, "rows[]", ret=True)
    msg(3.58, 3, 2, "prerequisitesMet: bool", ret=True)
    msg(3.92, 2, 3, "checkCapacity(courseId)", color=C["orange"])
    msg(4.22, 3, 4, "SELECT count(approved) for course", color=C["mid"])
    msg(4.52, 4, 3, "count", ret=True)
    msg(4.82, 3, 2, "hasCapacity: bool", ret=True)
    msg(5.15, 2, 3, "insertEnrollment(userId, courseId, PENDING)", color=C["green"])
    msg(5.45, 3, 4, "INSERT INTO enrollments…", color=C["mid"])
    msg(5.75, 4, 3, "enrollment row", ret=True)
    msg(6.08, 3, 2, "EnrollmentResponse", ret=True)
    msg(6.38, 2, 1, "EnrollmentResponse", ret=True)
    msg(6.68, 1, 0, "201 Created + enrollment body", ret=True)


# ── Experiment 5 ─────────────────────────────────────────────────────────────

def s_exp5(prs):
    _divider(prs, "Experiment 5", "State Chart Diagram",
             "Enrollment lifecycle state machine with guards and transitions")

    slide = _blank(prs)
    _header(slide, "Experiment 5 — State Chart Diagram",
            "Enrollment status: states, transitions, guards, and actions")
    _footer(slide)

    sw = Inches(2.2); sh = Inches(0.64)

    def state(l, t, label, fill, border):
        shp = _rect(slide, l, t, sw, sh, fill=fill, border=border, bw=Pt(2.2), radius=True)
        _stext(shp, label, size=Pt(13), bold=True, color=border)

    # Initial
    _ellipse(slide, Inches(4.65), Inches(1.2), Inches(0.32), Inches(0.32),
             fill=C["navy"], border=C["navy"])

    # PENDING
    state(Inches(3.9), Inches(1.72), "PENDING", C["lyell"], C["yellow"])
    # APPROVED
    state(Inches(0.85), Inches(3.8), "APPROVED", C["lgrn"], C["green"])
    # REJECTED
    state(Inches(6.95), Inches(3.8), "REJECTED", C["lorn"], C["orange"])
    # CANCELLED
    state(Inches(3.9), Inches(5.7), "CANCELLED", C["light"], C["mid"])

    # Final
    _ellipse(slide, Inches(4.65), Inches(6.6), Inches(0.32), Inches(0.32),
             fill=C["navy"], border=C["navy"])
    _ellipse(slide, Inches(4.59), Inches(6.54), Inches(0.44), Inches(0.44),
             fill=None, border=C["navy"], bw=Pt(2))

    # Transitions
    def tr(x1, y1, x2, y2, label, col=C["navy"]):
        c = slide.shapes.add_connector(1, x1, y1, x2, y2)
        c.line.color.rgb = col; c.line.width = Pt(1.5)
        mx, my = (x1+x2)/2, (y1+y2)/2
        _text(slide, mx-Inches(1.0), my-Inches(0.24), Inches(2.0), Inches(0.28),
              label, size=Pt(9.5), color=col, align=PP_ALIGN.CENTER)

    # init → PENDING
    _line(slide, Inches(4.81), Inches(1.52), Inches(4.81), Inches(1.72), color=C["navy"])

    # PENDING → APPROVED
    tr(Inches(3.9), Inches(1.72)+sh/2,
       Inches(0.85)+sw, Inches(3.8)+sh/2,
       "[capacity > 0][prereqs met]\nAdmin: approve", col=C["green"])

    # PENDING → REJECTED
    tr(Inches(3.9)+sw, Inches(1.72)+sh/2,
       Inches(6.95), Inches(3.8)+sh/2,
       "Admin: reject", col=C["orange"])

    # PENDING → CANCELLED (student cancel)
    _line(slide, Inches(3.9)+sw/2, Inches(1.72)+sh,
          Inches(3.9)+sw/2, Inches(5.7), color=C["mid"], w=Pt(1.4))
    _text(slide, Inches(5.15), Inches(3.55), Inches(1.5), Inches(0.28),
          "Student: cancel", size=Pt(9.5), italic=True, color=C["mid"])

    # APPROVED → CANCELLED
    tr(Inches(0.85)+sw/2, Inches(3.8)+sh,
       Inches(3.9)+sw/2-Inches(0.2), Inches(5.7),
       "Student: cancel\n[status=approved]", col=C["mid"])

    # APPROVED → Final (course completed, conceptual)
    _line(slide, Inches(0.85)+sw/2, Inches(3.8)+sh,
          Inches(4.81), Inches(6.6), color=C["green"], w=Pt(1))
    _text(slide, Inches(2.0), Inches(5.55), Inches(1.5), Inches(0.28),
          "Course ends", size=Pt(9), italic=True, color=C["green"])

    # REJECTED / CANCELLED → Final
    _line(slide, Inches(6.95)+sw/2, Inches(3.8)+sh,
          Inches(4.81), Inches(6.6), color=C["orange"], w=Pt(1))
    _line(slide, Inches(3.9)+sw/2, Inches(5.7)+sh,
          Inches(4.81), Inches(6.6), color=C["mid"], w=Pt(1))

    # Guard / action notes
    _rect(slide, Inches(0.3), Inches(6.7), Inches(9.4), Inches(0.42),
          fill=C["sky"], border=C["blue"], bw=Pt(1), radius=True)
    _text(slide, Inches(0.5), Inches(6.74), Inches(9.0), Inches(0.35),
          "Guards: [capacity > approved_count]  and  [all prerequisite courses have APPROVED enrollment]  "
          "are checked server-side at POST /enrollments and PUT /enrollments/:id/approve",
          size=Pt(10), color=C["navy"], wrap=True, va=MSO_ANCHOR.MIDDLE)


# ── Experiment 6 ─────────────────────────────────────────────────────────────

def s_exp6(prs):
    _divider(prs, "Experiment 6",
             "Layered Architecture + UML Package Diagram",
             "UI · Domain · Technical Services — partial logical architecture")

    slide = _blank(prs)
    _header(slide, "Experiment 6 — UML Package Diagram (Layered Architecture)",
            "Three-tier logical view with UML package notation")
    _footer(slide)

    layer_defs = [
        ("Presentation Layer  (User Interface)", C["blue"],  C["sky"],
         ["«component»\nReact SPA\n(client / client-b / client-c)",
          "«component»\nReact Router\n(role-aware routes)",
          "«component»\nAuth Context\n+ Route Guards",
          "«component»\nFeature Pages\n(Courses, Enrollments,\nAccount, Admin)",
          "«component»\nNotificationMenu\n+ Sonner toasts",
          "«component»\nAPI Client\n(lib/api.ts)"]),
        ("Domain Layer  (Business Logic)", C["teal"], C["lteal"],
         ["«service»\nCourseService\n(list, detail, CRUD)",
          "«service»\nEnrollmentService\n(request, approve/reject,\ncapacity + prereq checks)",
          "«service»\nStudentService\n(profile, scoped routes)",
          "«service»\nNotificationService\n(create, read)",
          "«service»\nAdminService\n(dashboard, reports)",
          "«entity»\nUser · Course · Enrollment\nNotification · RevokedToken"]),
        ("Technical Services Layer", C["navy"], C["light"],
         ["«service»\nAuthService\n(JWT sign/verify,\nbcrypt hash/verify)",
          "«service»\nDrizzle ORM\n(DB access layer)",
          "«infra»\nFastify Server\n(routes, guards,\nZod validation)",
          "«infra»\nPino Logger\n(structured logging)",
          "«infra»\nNeon PostgreSQL\n(pooled connection)",
          "«infra»\nUploadService\n(file validation +\nstorage helpers)"]),
    ]

    layer_h = Inches(1.95)
    for i, (title, hcol, bcol, comps) in enumerate(layer_defs):
        y = Inches(1.08) + i * (layer_h + Inches(0.12))
        # Package tab
        _rect(slide, Inches(0.3), y, Inches(2.5), Inches(0.35), fill=hcol)
        _text(slide, Inches(0.4), y+Inches(0.02), Inches(2.4), Inches(0.32),
              title, size=Pt(10.5), bold=True, color=C["white"], va=MSO_ANCHOR.MIDDLE)
        # Package body
        _rect(slide, Inches(0.3), y+Inches(0.35), Inches(9.4), layer_h-Inches(0.35),
              fill=bcol, border=hcol, bw=Pt(1.8))
        # Components inside
        cw = Inches(1.47); ch = Inches(1.28)
        for j, cname in enumerate(comps):
            cx = Inches(0.5) + j * (cw + Inches(0.06))
            cy = y + Inches(0.45)
            _boxt(slide, cx, cy, cw, ch, cname, fill=C["white"],
                  border=hcol, bw=Pt(1.2), size=Pt(9), tc=C["dark"])

    # Dependency arrows between layers
    mid_x = Inches(5.0)
    for yi in [Inches(1.08)+layer_h, Inches(1.08)+2*(layer_h+Inches(0.12))]:
        c = slide.shapes.add_connector(1, mid_x, yi+Inches(0.12),
                                       mid_x, yi+Inches(0.12)+Inches(0.1))
        c.line.color.rgb = C["navy"]; c.line.width = Pt(2)
    _text(slide, Inches(5.5), Inches(3.05), Inches(2.0), Inches(0.28),
          "«uses»", size=Pt(10), italic=True, color=C["navy"])
    _text(slide, Inches(5.5), Inches(4.98), Inches(2.0), Inches(0.28),
          "«uses»", size=Pt(10), italic=True, color=C["navy"])


# ── Experiment 7 ─────────────────────────────────────────────────────────────

def s_exp7(prs):
    _divider(prs, "Experiment 7",
             "Technical Services & Domain Objects Layer",
             "Implementation details for the backend API and database access")

    slide = _blank(prs)
    _header(slide, "Experiment 7 — Technical Services & Domain Layer Implementation",
            "Fastify routes · Drizzle ORM · services · domain entities")
    _footer(slide)

    _card(slide, Inches(0.3), Inches(1.08), Inches(3.0), Inches(6.0),
          "Domain Entities (Drizzle Schema)",
          ["users\n  id, name, email, passwordHash,\n  role, isActive",
           "courses\n  id, code, title, description,\n  credits, capacity",
           "course_prerequisites\n  courseId (FK), prerequisiteId (FK)",
           "enrollments\n  id, userId (FK), courseId (FK),\n  status, createdAt",
           "notifications\n  id, userId (FK), message,\n  type, read, createdAt",
           "revoked_tokens\n  id, jti, expiresAt"],
          hfill=C["navy"])

    _card(slide, Inches(3.5), Inches(1.08), Inches(3.1), Inches(6.0),
          "Technical Services",
          ["AuthService\n  bcrypt.hash / compare\n  jose sign / verify JWT\n  jti revocation check",
           "DrizzleORM\n  Type-safe query builder\n  Migrations via drizzle-kit\n  Neon pooled driver",
           "Fastify preHandlers\n  authenticate(req, reply)\n  requireRole('admin'|'student')",
           "Zod Schemas\n  registerSchema, loginSchema\n  enrollmentSchema, courseSchema\n  Validated on every route",
           "Pino Logger\n  LOG_LEVEL env var\n  Redacts Authorization header\n  pino-pretty in dev"],
          hfill=C["blue"])

    _card(slide, Inches(6.8), Inches(1.08), Inches(2.85), Inches(6.0),
          "Key Implementation Details",
          ["JWT payload: { sub, role, jti, exp }",
           "Logout: INSERT jti into revoked_tokens",
           "Approval guard:\n  SELECT COUNT(*) approved < capacity",
           "Prereq guard:\n  check APPROVED enrollments",
           "Notification in same DB\n  transaction as status update",
           "UNIQUE(userId, courseId)\n  prevents duplicate enrollments",
           "GET /enrollments/mine\n  scoped to authenticated user",
           "GET /admin/reports/*\n  aggregated SQL queries",
           "File upload:\n  PDF validate + write to FS"],
          hfill=C["teal"])


# ── Experiment 8 ─────────────────────────────────────────────────────────────

def s_exp8(prs):
    _divider(prs, "Experiment 8", "User Interface Layer Implementation",
             "React SPA · routing · pages · components · state management")

    slide = _blank(prs)
    _header(slide, "Experiment 8 — User Interface Layer Implementation",
            "React + Vite · shadcn/ui (Base UI) · Tailwind · role-aware routing")
    _footer(slide)

    # Route tree diagram
    _rect(slide, Inches(0.3), Inches(1.08), Inches(4.55), Inches(0.38),
          fill=C["navy"])
    _text(slide, Inches(0.4), Inches(1.1), Inches(4.35), Inches(0.35),
          "Student Route Tree", size=Pt(12), bold=True, color=C["white"],
          va=MSO_ANCHOR.MIDDLE)
    routes_s = [
        ("/               → Home / Landing"),
        ("/courses        → Course catalog list"),
        ("/courses/:id    → Course detail + Enroll CTA"),
        ("/enrollments    → My enrollments (polls 15s)"),
        ("/account        → Profile edit + file upload"),
        ("/login  /register → Auth forms"),
        ("NotificationMenu → Bell (all authenticated pages)"),
    ]
    for i, r in enumerate(routes_s):
        bg = C["sky"] if i % 2 == 0 else C["white"]
        _rect(slide, Inches(0.3), Inches(1.46)+i*Inches(0.38), Inches(4.55),
              Inches(0.38), fill=bg, border=C["border"], bw=Pt(0.4))
        _text(slide, Inches(0.4), Inches(1.48)+i*Inches(0.38), Inches(4.35),
              Inches(0.35), r, size=Pt(11), color=C["dark"], va=MSO_ANCHOR.MIDDLE)

    _rect(slide, Inches(5.05), Inches(1.08), Inches(4.6), Inches(0.38),
          fill=C["teal"])
    _text(slide, Inches(5.15), Inches(1.1), Inches(4.4), Inches(0.35),
          "Admin Route Tree (/admin/*)", size=Pt(12), bold=True, color=C["white"],
          va=MSO_ANCHOR.MIDDLE)
    routes_a = [
        ("/admin/dashboard     → KPIs + send notification"),
        ("/admin/students      → Student list + detail"),
        ("/admin/courses       → Course CRUD"),
        ("/admin/enrollments   → Approve / reject queue"),
        ("/admin/reports       → Enrollment, student, course"),
        ("/admin/admins        → Admin accounts management"),
        ("AdminLayout wraps all /admin/* routes (RBAC)"),
    ]
    for i, r in enumerate(routes_a):
        bg = C["lteal"] if i % 2 == 0 else C["white"]
        _rect(slide, Inches(5.05), Inches(1.46)+i*Inches(0.38), Inches(4.6),
              Inches(0.38), fill=bg, border=C["border"], bw=Pt(0.4))
        _text(slide, Inches(5.15), Inches(1.48)+i*Inches(0.38), Inches(4.4),
              Inches(0.35), r, size=Pt(11), color=C["dark"], va=MSO_ANCHOR.MIDDLE)

    # Component library
    _card(slide, Inches(0.3), Inches(4.48), Inches(4.55), Inches(2.6),
          "shadcn/ui Component Library (Base UI)",
          ["Button, Input, Card, Dialog, Dropdown — from shadcn CLI",
           "Base UI primitives (not Radix) — composition pattern",
           "Tailwind CSS theming via CSS variables",
           "lucide-react icons throughout",
           "Sonner toasts for async feedback (enroll, approve, errors)",
           "All components in src/components/ui/ (generated)"],
          hfill=C["blue"])

    _card(slide, Inches(5.05), Inches(4.48), Inches(4.6), Inches(2.6),
          "State & Data Fetching",
          ["Auth state in React Context (AuthContext)",
           "JWT stored in-memory / localStorage (lab scope)",
           "API calls via fetch in lib/api.ts (no TanStack Query)",
           "Bearer token attached to every protected request",
           "Enrollment list polls via setInterval (15 s) while visible",
           "React Router v6 for client-side navigation + guards"],
          hfill=C["purple"])


# ── Experiment 9 ─────────────────────────────────────────────────────────────

def s_exp9(prs):
    _divider(prs, "Experiment 9",
             "Component & Deployment Diagrams",
             "Physical components and deployment topology")

    slide = _blank(prs)
    _header(slide, "Experiment 9 — Component Diagram",
            "System components, interfaces, and dependencies")
    _footer(slide)

    # Frontend components box
    _rect(slide, Inches(0.3), Inches(1.08), Inches(4.3), Inches(0.4), fill=C["blue"])
    _text(slide, Inches(0.4), Inches(1.1), Inches(4.1), Inches(0.37),
          "«subsystem»  Frontend (Vite/React)", size=Pt(12), bold=True, color=C["white"],
          va=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(0.3), Inches(1.48), Inches(4.3), Inches(4.3),
          fill=C["sky"], border=C["blue"], bw=Pt(1.5))
    fe_comps = ["AppRouter + Layouts", "AuthContext / Guards",
                "CoursesPage + EnrollmentsPage",
                "AccountPage + AdminPages",
                "NotificationMenu (bell)",
                "lib/api.ts (HTTP client)"]
    for i, c in enumerate(fe_comps):
        _boxt(slide, Inches(0.45), Inches(1.6)+i*Inches(0.6), Inches(4.0),
              Inches(0.5), c, fill=C["white"], border=C["blue"],
              bw=Pt(1), size=Pt(11.5), tc=C["navy"])

    # Arrow interface
    _line(slide, Inches(4.6), Inches(3.3), Inches(5.15), Inches(3.3),
          color=C["navy"], w=Pt(2))
    _text(slide, Inches(4.52), Inches(3.0), Inches(0.72), Inches(0.28),
          "REST\nHTTPS+JWT", size=Pt(9), italic=True, color=C["navy"],
          align=PP_ALIGN.CENTER)

    # Backend components box
    _rect(slide, Inches(5.15), Inches(1.08), Inches(4.5), Inches(0.4), fill=C["teal"])
    _text(slide, Inches(5.25), Inches(1.1), Inches(4.3), Inches(0.37),
          "«subsystem»  Backend (Fastify API)", size=Pt(12), bold=True, color=C["white"],
          va=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(5.15), Inches(1.48), Inches(4.5), Inches(4.3),
          fill=C["lteal"], border=C["teal"], bw=Pt(1.5))
    be_comps = ["Route Handlers (auth, courses,\nenrollments, students, admin)",
                "preHandlers: authenticate\n+ requireRole",
                "EnrollmentService\n(capacity + prereq logic)",
                "DrizzleORM (type-safe queries)",
                "AuthService (JWT + bcrypt)",
                "UploadService (PDF storage)"]
    for i, c in enumerate(be_comps):
        _boxt(slide, Inches(5.3), Inches(1.6)+i*Inches(0.6), Inches(4.15),
              Inches(0.5), c, fill=C["white"], border=C["teal"],
              bw=Pt(1), size=Pt(11), tc=C["dark"])

    # DB
    _line(slide, Inches(7.4), Inches(5.78), Inches(7.4), Inches(6.12),
          color=C["navy"], w=Pt(2))
    _boxt(slide, Inches(5.15), Inches(6.12), Inches(4.5), Inches(0.58),
          "«database»  PostgreSQL on Neon (pooled)",
          fill=C["lyell"], border=C["yellow"], bw=Pt(1.5), size=Pt(12), tc=C["dark"])


def s_exp9_deploy(prs):
    slide = _blank(prs)
    _header(slide, "Experiment 9 — Deployment Diagram",
            "Runtime topology: browser · Vite dev servers · Fastify API · Neon DB")
    _footer(slide)

    def node(l, t, w, h, label, fill, border):
        _rect(slide, l, t, w, h, fill=fill, border=border, bw=Pt(2))
        # 3D tab
        _rect(slide, l, t, w, Inches(0.35), fill=border)
        _text(slide, l+Inches(0.1), t+Inches(0.02), w-Inches(0.18), Inches(0.32),
              label, size=Pt(11), bold=True, color=C["white"], va=MSO_ANCHOR.MIDDLE)

    # Client Browser
    node(Inches(3.8), Inches(1.08), Inches(2.4), Inches(1.25),
         "«device» User Browser", C["light"], C["navy"])
    _text(slide, Inches(3.95), Inches(1.52), Inches(2.15), Inches(0.65),
          "React SPA (HTML + JS)\nRendered in browser", size=Pt(10.5), color=C["dark"])

    # Three Vite servers
    for j, (port, lbl) in enumerate([("5173","client"),("5174","client-b"),("5175","client-c")]):
        x = Inches(0.3) + j*Inches(3.22)
        node(x, Inches(2.7), Inches(2.9), Inches(1.3),
             f"«server» Vite Dev — :{port}", C["sky"], C["blue"])
        _text(slide, x+Inches(0.15), Inches(3.12), Inches(2.65), Inches(0.7),
              f"React/Vite SPA\n{lbl}/  (static in prod)", size=Pt(10.5), color=C["dark"])
        _line(slide, Inches(3.8)+Inches(0.45), Inches(2.33),
              x+Inches(1.35), Inches(2.7), color=C["blue"], w=Pt(1.2))

    # Fastify API server
    node(Inches(3.4), Inches(4.4), Inches(3.2), Inches(1.35),
         "«server» Fastify API — :3000", C["lteal"], C["teal"])
    _text(slide, Inches(3.55), Inches(4.82), Inches(2.95), Inches(0.82),
          "Node.js 20+  ·  TypeScript\nRoutes, Services, Drizzle ORM\nPino logging · Zod validation",
          size=Pt(11), color=C["dark"])
    # arrows Vite → API
    for j in range(3):
        x = Inches(0.3) + j*Inches(3.22)
        _line(slide, x+Inches(1.35), Inches(4.0),
              Inches(5.0), Inches(4.4), color=C["teal"], w=Pt(1.2))
    _text(slide, Inches(1.5), Inches(4.1), Inches(2.0), Inches(0.26),
          "HTTPS + Bearer JWT", size=Pt(9.5), italic=True, color=C["teal"])

    # Neon PostgreSQL
    node(Inches(3.4), Inches(6.05), Inches(3.2), Inches(1.05),
         "«cloud DB» Neon PostgreSQL", C["lyell"], C["yellow"])
    _text(slide, Inches(3.55), Inches(6.45), Inches(2.95), Inches(0.55),
          "Managed serverless Postgres\nPooled connection URL (DATABASE_URL)",
          size=Pt(11), color=C["dark"])
    _line(slide, Inches(5.0), Inches(5.75), Inches(5.0), Inches(6.05),
          color=C["yellow"], w=Pt(1.8))
    _text(slide, Inches(5.1), Inches(5.82), Inches(1.5), Inches(0.25),
          "SQL via Drizzle", size=Pt(9.5), italic=True, color=C["yellow"])


# ── Experiment 10 ────────────────────────────────────────────────────────────

def s_exp10(prs):
    _divider(prs, "Experiment 10",
             "Testing, Documentation & Final Product",
             "Vitest · 80 manual test cases · CI · OpenAPI · README")

    slide = _blank(prs)
    _header(slide, "Experiment 10 — Testing & Quality Results",
            "Automated tests, manual test matrix, and CI pipeline")
    _footer(slide)

    # Test pyramid
    layers = [
        (Inches(5.0), Inches(0.58), "Unit Tests  (Vitest)", C["blue"], C["sky"],
         "password hashing · JWT sign/verify"),
        (Inches(4.0), Inches(0.58), "HTTP Smoke  (app.inject)", C["teal"], C["lteal"],
         "GET /health — no DB required"),
        (Inches(3.2), Inches(0.58), "DB Integration  (optional)", C["purple"], C["lpurp"],
         "register + login (RUN_DB_INTEGRATION=1)"),
        (Inches(2.2), Inches(0.55), "Manual / Exploratory  (80 TCs)", C["orange"], C["lorn"],
         "TESTING_REPORT.md — all Pass"),
    ]
    y = Inches(1.08)
    for w, h, label, hfill, bfill, sub in layers:
        x = (W - w) / 2
        _rect(slide, x, y, w, h, fill=bfill, border=hfill, bw=Pt(1.8))
        _text(slide, x+Inches(0.1), y+Inches(0.01), w-Inches(0.2), h-Inches(0.05),
              label, size=Pt(11.5), bold=True, color=hfill, align=PP_ALIGN.CENTER,
              va=MSO_ANCHOR.MIDDLE)
        _text(slide, x+w+Inches(0.15), y+Inches(0.12), Inches(3.5), h,
              sub, size=Pt(10), italic=True, color=C["mid"])
        y += h + Inches(0.06)

    # Test summary table
    _rect(slide, Inches(0.3), Inches(3.65), Inches(4.5), Inches(0.38), fill=C["navy"])
    for col, (lbl, off) in enumerate([("Metric",1.0),("Value",3.2)]):
        _text(slide, Inches(0.4+off-1.0), Inches(3.67), Inches(2.2), Inches(0.35),
              lbl, size=Pt(12), bold=True, color=C["white"], va=MSO_ANCHOR.MIDDLE)
    rows = [("Total test cases","80"),("Passed","80"),
            ("Failed","0"),("Pass rate","100 %"),
            ("Modules covered","12")]
    for i, (k, v) in enumerate(rows):
        bg = C["sky"] if i%2==0 else C["white"]
        y2 = Inches(4.03) + i*Inches(0.38)
        _rect(slide, Inches(0.3), y2, Inches(4.5), Inches(0.38),
              fill=bg, border=C["border"], bw=Pt(0.5))
        _text(slide, Inches(0.45), y2+Inches(0.02), Inches(2.5), Inches(0.35),
              k, size=Pt(12), color=C["dark"], va=MSO_ANCHOR.MIDDLE)
        _text(slide, Inches(3.0), y2+Inches(0.02), Inches(1.6), Inches(0.35),
              v, size=Pt(12), bold=True, color=C["navy"], va=MSO_ANCHOR.MIDDLE,
              align=PP_ALIGN.CENTER)

    # CI + Docs
    _card(slide, Inches(5.05), Inches(3.65), Inches(4.6), Inches(2.18),
          "CI Pipeline (.github/workflows/ci.yml)",
          ["Trigger: push or PR to main / master",
           "Step 1 — npm install (workspace)",
           "Step 2 — lint (workspaces that define it)",
           "Step 3 — npm test -w server (Vitest)",
           "Step 4 — npm run build (all workspaces)"],
          hfill=C["teal"])

    _card(slide, Inches(5.05), Inches(5.98), Inches(4.6), Inches(1.25),
          "Documentation",
          ["OpenAPI/Swagger at GET /documentation (auto-gen)",
           "README: setup, env vars, migrate, first-admin steps",
           "PROJECT_PLAN.md: full scope + phased roadmap"],
          hfill=C["blue"])


def s_exp10_final(prs):
    slide = _blank(prs)
    _header(slide, "Experiment 10 — Final Product Overview",
            "Delivered system: all phases complete, all tests pass")
    _footer(slide)

    phases_done = [
        ("Phase 1  Foundation",
         "Workspaces, Neon + users migration, GET /health, Vite + shadcn (Base UI)"),
        ("Phase 2  Authentication",
         "POST /auth/register, login, GET /auth/me, logout (jti revocation), change-password, auth UI"),
        ("Phase 3  Courses",
         "courses + course_prerequisites tables, CRUD routes, public catalog, admin course UI"),
        ("Phase 4  Enrollments",
         "POST /enrollments with capacity + prereq + duplicate checks; approve/reject/cancel; status UI"),
        ("Phase 5  Admin Aggregation",
         "GET /students, /admin/dashboard, /admin/students|courses|enrollments|reports|admins"),
        ("Phase 6  Notifications",
         "Notification table; approve/reject creates notification in same DB transaction; bell UI + toasts"),
        ("Phase 7  Quality",
         "Vitest tests, optional DB integration, OpenAPI/Swagger, GitHub Actions CI, README"),
    ]
    for i, (phase, desc) in enumerate(phases_done):
        y = Inches(1.1) + i * Inches(0.84)
        _rect(slide, Inches(0.3), y, Inches(0.62), Inches(0.62),
              fill=C["green"], border=None, radius=True)
        _text(slide, Inches(0.3), y+Inches(0.06), Inches(0.62), Inches(0.5),
              "✓", size=Pt(20), bold=True, color=C["white"],
              align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
        _rect(slide, Inches(1.05), y, Inches(8.65), Inches(0.62),
              fill=C["lgrn"] if i%2==0 else C["sky"],
              border=C["border"], bw=Pt(0.6))
        _text(slide, Inches(1.18), y+Inches(0.0), Inches(2.5), Inches(0.32),
              phase, size=Pt(12), bold=True, color=C["navy"], va=MSO_ANCHOR.MIDDLE)
        _text(slide, Inches(1.18), y+Inches(0.3), Inches(8.42), Inches(0.32),
              desc, size=Pt(11), color=C["dark"])

    _rect(slide, 0, Inches(7.0), W, Inches(0.5), fill=C["teal"])
    _text(slide, Inches(0.4), Inches(7.04), Inches(9.2), Inches(0.44),
          "All 10 experiments completed  ·  100% test pass rate  ·  Full-stack, production-ready architecture",
          size=Pt(15), bold=True, color=C["white"], align=PP_ALIGN.CENTER,
          va=MSO_ANCHOR.MIDDLE)


# ─────────────────────────── main ────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    s_title(prs)
    s_intro(prs)
    s_objective(prs)

    s_exp1(prs)           # Exp 1: Problem statement

    s_exp2_srs(prs)       # Exp 2: IEEE SRS
    s_exp2_risk_gantt(prs)

    s_exp3_usecase(prs)   # Exp 3: Use Case + Activity
    s_exp3_activity(prs)

    s_exp4_class(prs)     # Exp 4: Class + Interaction
    s_exp4_seq(prs)

    s_exp5(prs)           # Exp 5: State Chart

    s_exp6(prs)           # Exp 6: Package diagram

    s_exp7(prs)           # Exp 7: Technical + Domain impl

    s_exp8(prs)           # Exp 8: UI layer

    s_exp9(prs)           # Exp 9: Component + Deployment
    s_exp9_deploy(prs)

    s_exp10(prs)          # Exp 10: Testing + docs
    s_exp10_final(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT.resolve())


if __name__ == "__main__":
    main()
